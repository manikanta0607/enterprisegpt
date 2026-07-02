"""PPTX text extraction using python-pptx."""

import io

from pptx import Presentation

from app.core.exceptions import ValidationError
from app.services.parsers.base import DocumentParser


class PptxParser(DocumentParser):
    """Extracts text from PowerPoint (.pptx) files, slide by slide."""

    def extract_text(self, data: bytes) -> str:
        """Extract text from all text-bearing shapes across all slides.

        Args:
            data: The raw .pptx file bytes.

        Returns:
            Extracted text, with each slide's content separated by a double
            newline and prefixed with a slide marker.

        Raises:
            ValidationError: If the file cannot be parsed as a valid .pptx.
        """
        try:
            presentation = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ValidationError("File could not be read as a valid PPTX presentation") from exc

        slides_text = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            lines = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if lines:
                slides_text.append(f"[Slide {slide_number}]\n" + "\n".join(lines))

        return "\n\n".join(slides_text).strip()
