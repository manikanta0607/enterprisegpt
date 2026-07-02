"""Tests for document parsers.

Fixture files are generated in-memory (not checked into the repo as binary
files) using the same libraries a real user's documents would come from,
so these tests exercise the actual parsing libraries end-to-end.
"""

import io

import pytest
from docx import Document as DocxDocument
from fpdf import FPDF
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from app.core.exceptions import ValidationError
from app.services.parsers.docx_parser import DocxParser
from app.services.parsers.factory import get_parser_for_content_type
from app.services.parsers.ocr_parser import OcrParser
from app.services.parsers.pdf_parser import PdfParser
from app.services.parsers.pptx_parser import PptxParser


def _make_pdf_bytes(text: str) -> bytes:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, text)
    return bytes(pdf.output())


def _make_docx_bytes(text: str) -> bytes:
    doc = DocxDocument()
    doc.add_paragraph(text)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _make_pptx_bytes(text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _make_image_bytes(text: str) -> bytes:
    image = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((10, 40), text, fill="black")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def test_pdf_parser_extracts_text():
    """PdfParser should extract text from a real generated PDF."""
    data = _make_pdf_bytes("Hello from EnterpriseGPT PDF test")
    result = PdfParser().extract_text(data)
    assert "Hello from EnterpriseGPT" in result


def test_pdf_parser_rejects_invalid_file():
    """PdfParser should raise ValidationError for non-PDF bytes."""
    with pytest.raises(ValidationError):
        PdfParser().extract_text(b"not a real pdf")


def test_docx_parser_extracts_text():
    """DocxParser should extract paragraph text from a real .docx file."""
    data = _make_docx_bytes("Hello from EnterpriseGPT DOCX test")
    result = DocxParser().extract_text(data)
    assert "Hello from EnterpriseGPT DOCX test" in result


def test_docx_parser_rejects_invalid_file():
    """DocxParser should raise ValidationError for non-DOCX bytes."""
    with pytest.raises(ValidationError):
        DocxParser().extract_text(b"not a real docx")


def test_pptx_parser_extracts_text():
    """PptxParser should extract textbox content from a real .pptx file."""
    data = _make_pptx_bytes("Hello from EnterpriseGPT PPTX test")
    result = PptxParser().extract_text(data)
    assert "Hello from EnterpriseGPT PPTX test" in result
    assert "[Slide 1]" in result


def test_pptx_parser_rejects_invalid_file():
    """PptxParser should raise ValidationError for non-PPTX bytes."""
    with pytest.raises(ValidationError):
        PptxParser().extract_text(b"not a real pptx")


def test_ocr_parser_extracts_text_from_image():
    """OcrParser should recognize rendered text in a real PNG via Tesseract."""
    data = _make_image_bytes("HELLO OCR")
    result = OcrParser().extract_text(data)
    assert "HELLO" in result.upper()


def test_ocr_parser_rejects_invalid_file():
    """OcrParser should raise ValidationError for non-image bytes."""
    with pytest.raises(ValidationError):
        OcrParser().extract_text(b"not a real image")


def test_factory_returns_correct_parser_per_content_type():
    """The parser factory should route each MIME type to the right parser."""
    assert isinstance(get_parser_for_content_type("application/pdf"), PdfParser)
    assert isinstance(
        get_parser_for_content_type(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
        DocxParser,
    )
    assert isinstance(
        get_parser_for_content_type(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        PptxParser,
    )
    assert isinstance(get_parser_for_content_type("image/png"), OcrParser)


def test_factory_rejects_unsupported_content_type():
    """The parser factory should raise for an unregistered MIME type."""
    with pytest.raises(ValidationError):
        get_parser_for_content_type("application/zip")
